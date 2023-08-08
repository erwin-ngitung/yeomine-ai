from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP
from datetime import datetime as dt
from streamlit import session_state as state
import pytz


def model_analysis(img_path, ppt_template):
    time_now = pytz.timezone('Asia/Jakarta')
    true_time = dt.now(time_now).strftime('%d-%m-%Y')

    prs = Presentation(ppt_template)

    slide_0 = prs.slides[0]

    # Slide 0
    line_0 = slide_0.shapes[2].text_frame.paragraphs[0]
    line_0.text = f'Date: {true_time}'
    line_0.alignment = PP.CENTER
    line_0.font.name = 'Archive Black'
    line_0.font.size = Pt(35)
    line_0.font.bold = True

    idx_slide = [2, 3, 4, 5, 6, 7, 9, 10, 11, 12, 13, 14]
    name_slide = {'confusion_matrix_normalized': 'Confusion Matrix',
                  'F1_curve': 'F1 Curve',
                  'P_curve': 'P Curve',
                  'PR_curve': 'PR Curve',
                  'R_curve': 'R Curve',
                  'results': 'Summary',
                  'labels': 'Label Validation',
                  'train_batch0': 'Train Batch 0 Validation',
                  'train_batch1': 'Train Batch 1 Validation',
                  'train_batch2': 'Train Batch 2 Validation',
                  'val_batch0_labels': 'Test Batch 0 Validation',
                  'val_batch0_pred': 'Prediction Batch 0 Validation'}

    for i, idx in enumerate(idx_slide):
        slides = prs.slides[idx]
        line = slides.shapes[2].text_frame.paragraphs[0]
        line.text = list(name_slide.values())[i]
        line.alignment = PP.LEFT
        line.font.name = 'Archive Black'
        line.font.size = Pt(50)
        line.font.bold = True

        path_img = f'{img_path}/{list(name_slide.keys())[i]}.png'

        picture = slides.shapes
        picture.add_picture(path_img,
                            left=Inches(2.5),
                            top=Inches(2),
                            width=Inches(15),
                            height=Inches(9))

    return prs


def report_analysis(img_path, ppt_template):
    time_now = pytz.timezone('Asia/Jakarta')
    true_time = dt.now(time_now).strftime('%d-%m-%Y')

    prs = Presentation(ppt_template)

    slide_0 = prs.slides[0]

    # Slide 0
    line_0 = slide_0.shapes[2].text_frame.paragraphs[0]
    line_0.text = f'Date: {true_time}'
    line_0.alignment = PP.CENTER
    line_0.font.name = 'Archive Black'
    line_0.font.size = Pt(35)
    line_0.font.bold = True

    idx_slide = [1, 2]
    name_slide = {'coordinate-object': 'Graph Coordinate Object',
                  'graph-count-object': 'Graph Count Object'}

    for i, idx in enumerate(idx_slide):
        slides = prs.slides[idx]
        line = slides.shapes[2].text_frame.paragraphs[0]
        line.text = list(name_slide.values())[i]
        line.alignment = PP.LEFT
        line.font.name = 'Archive Black'
        line.font.size = Pt(50)
        line.font.bold = True

        path_img = f'{img_path}/{list(name_slide.keys())[i]}'

        picture = slides.shapes
        picture.add_picture(path_img,
                            left=Inches(2.5),
                            top=Inches(2),
                            width=Inches(15),
                            height=Inches(9))

    return prs
